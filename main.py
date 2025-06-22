import os
import faiss
import numpy as np
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from langchain.text_splitter import RecursiveCharacterTextSplitter
import pdfplumber
import DouBao
# ---------- 文档解析 ----------
def extract_text(path):
    ext = os.path.splitext(path)[-1].lower()
    if ext == ".docx":
        doc = Document(path)
        return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext == ".pptx":
        prs = Presentation(path)
        return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    elif ext == ".pdf":
        text = ""
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n"
        return text
    elif ext == ".txt":
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    else:
        raise ValueError(f"不支持的文件类型: {ext}")

# ---------- 向量生成 ----------
def embed_text_chunks(text_chunks, model_name='all-MiniLM-L6-v2', max_retries=3):
    from time import sleep
    for attempt in range(max_retries):
        try:
            model = SentenceTransformer(model_name)
            embeddings = model.encode(text_chunks)
            return embeddings, model
        except Exception as e:
            if attempt < max_retries - 1:
                wait_time = 2 ** attempt  # Exponential backoff
                print(f"尝试加载模型失败，将在 {wait_time} 秒后重试... ({attempt+1}/{max_retries})")
                sleep(wait_time)
            else:
                print(f"无法加载模型，请检查网络连接。错误: {str(e)}")
                raise

# ---------- 索引构建 ----------
def build_index(vectors, save_path="index.faiss"):
    index = faiss.IndexFlatL2(vectors.shape[1])
    index.add(vectors)
    faiss.write_index(index, save_path)

# ---------- 执行搜索 ----------
def search_chunks(query, model, chunks, index_path="index.faiss"):
    index = faiss.read_index(index_path)
    query_vec = model.encode([query])
    D, I = index.search(np.array(query_vec), k=3)
    return [chunks[i] for i in I[0]]

# ---------- 调用 GPT 模型回答 ----------
def generate_answer(context_chunks, query):
    context = "\n".join(context_chunks)
    prompt = f"""你是一个教学助手。根据以下资料内容回答问题：
资料：
{context}

问题：{query}

请用简洁清晰的语言回答："""
    # print(f"外部知识库检索到的相关内容：\n{context}\n")
    res = DouBao.get_answer(prompt)
    return res

# ---------- 主流程：构建知识库 ----------
def build_knowledge(status_callback=None):
    all_text = ""
    if not os.path.exists("docs"):
        os.makedirs("docs")
        if status_callback:
            status_callback("创建了docs文件夹，请将文档放入后重试")
        return

    files = os.listdir("docs")
    if not files:
        if status_callback:
            status_callback("docs文件夹为空，请添加文档后重试")
        return

    for filename in files:
        path = os.path.join("docs", filename)
        try:
            text = extract_text(path)
            all_text += text + "\n"
            if status_callback:
                status_callback(f"正在解析：{filename}")
            else:
                print(f"正在解析：{filename}")
        except Exception as e:
            msg = f"跳过 {filename}：{e}"
            if status_callback:
                status_callback(msg)
            else:
                print(msg)

    if status_callback:
        status_callback("正在分割文本...")
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_text(all_text)

    if status_callback:
        status_callback("正在向量化文本并构建索引...")
    vectors, model = embed_text_chunks(chunks)
    build_index(np.array(vectors))

    with open("chunks.txt", "w", encoding="utf-8") as f:
        f.write("\n\n===CHUNK===\n\n".join(chunks))

    complete_msg = "✅ 构建完成！已保存索引与文本块。"
    if status_callback:
        status_callback(complete_msg)
    else:
        print(complete_msg)

# ---------- 主流程：提问查询 ----------
def ask_question(query, callback=None):
    if not os.path.exists("chunks.txt") or not os.path.exists("index.faiss"):
        result = "请先构建知识库！"
        if callback:
            callback(result)
            return
        else:
            return result

    with open("chunks.txt", "r", encoding="utf-8") as f:
        chunks = f.read().split("\n\n===CHUNK===\n\n")
    from time import sleep
    for attempt in range(3):  # Try 3 times
        try:
            model = SentenceTransformer('all-MiniLM-L6-v2')
            break
        except Exception as e:
            if attempt < 2:
                wait_time = 2 ** attempt
                print(f"加载模型失败，将在 {wait_time} 秒后重试... ({attempt + 1}/3)")
                sleep(wait_time)
            else:
                if callback:
                    callback(f"无法加载模型，请检查网络连接。错误: {str(e)}")
                    return
                else:
                    return f"无法加载模型，请检查网络连接。错误: {str(e)}"

    results = search_chunks(query, model, chunks)
    answer = generate_answer(results, query)

    formatted_answer = f"\n【问题】：{query}\n【回答】：\n{answer}"
    if callback:
        callback(formatted_answer)
    else:
        print(formatted_answer)
        return answer

# ---------- 命令行入口 ----------
def cli_interface():
    while True:
        print("\n请选择操作：")
        print("1. 构建知识库")
        print("2. 提问查询")
        print("3. 退出")
        choice = input("请输入数字选择操作：")

        if choice == "1":
            build_knowledge()
        elif choice == "2":
            query = input("请输入你的问题：")
            ask_question(query)
        elif choice == "3":
            print("已退出程序。")
            break
        else:
            print("无效选择，请重新输入。")

# ---------- GUI 界面 ----------
def gui_interface():
    import tkinter as tk
    from tkinter import scrolledtext, messagebox, filedialog, ttk
    import threading

    def update_status(message):
        status_text.config(state=tk.NORMAL)
        status_text.insert(tk.END, message + "\n")
        status_text.see(tk.END)
        status_text.config(state=tk.DISABLED)

    def display_answer(answer):
        answer_text.config(state=tk.NORMAL)
        answer_text.delete(1.0, tk.END)
        answer_text.insert(tk.END, answer)
        answer_text.config(state=tk.DISABLED)

    def on_build_click():
        status_text.config(state=tk.NORMAL)
        status_text.delete(1.0, tk.END)
        status_text.config(state=tk.DISABLED)

        build_button.config(state=tk.DISABLED)
        progress_bar.grid(row=0, column=0, columnspan=2, padx=5, pady=5, sticky='ew')
        progress_bar.start()

        def build_thread():
            try:
                def progress_callback(message, progress=None):
                    update_status(message)
                    if progress is not None and isinstance(progress, float):
                        if 0 <= progress <= 1:
                            # Convert to percentage for better display
                            percentage = int(progress * 100)
                            progress_var.set(percentage)
                            progress_label.config(text=f"{percentage}%")

                # Modified build_knowledge call with progress reporting
                build_knowledge_with_progress(progress_callback)
            except Exception as e:
                update_status(f"错误: {str(e)}")
            finally:
                build_button.config(state=tk.NORMAL)
                progress_bar.stop()
                progress_bar.grid_remove()

        threading.Thread(target=build_thread).start()

    def build_knowledge_with_progress(callback):
        """Enhanced version of build_knowledge with progress reporting"""
        all_text = ""
        if not os.path.exists("docs"):
            os.makedirs("docs")
            callback("创建了docs文件夹，请将文档放入后重试")
            return

        files = os.listdir("docs")
        if not files:
            callback("docs文件夹为空，请添加文档后重试")
            return

        # Parse documents
        total_files = len(files)
        for i, filename in enumerate(files):
            path = os.path.join("docs", filename)
            try:
                text = extract_text(path)
                all_text += text + "\n"
                callback(f"正在解析：{filename}", progress=(i / total_files) * 0.3)
            except Exception as e:
                callback(f"跳过 {filename}：{e}")

        callback("正在分割文本...", progress=0.35)
        splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        chunks = splitter.split_text(all_text)
        callback(f"文本已分割为 {len(chunks)} 个块", progress=0.4)

        callback("正在向量化文本...", progress=0.45)

        # Custom embed function to report progress
        def embed_with_progress(text_chunks, model_name='all-MiniLM-L6-v2', max_retries=3):
            from time import sleep
            for attempt in range(max_retries):
                try:
                    callback("加载向量模型中...", progress=0.5)
                    model = SentenceTransformer(model_name)

                    # Process chunks in batches to report progress
                    total_chunks = len(text_chunks)
                    batch_size = 50  # Adjust based on your needs
                    all_embeddings = []

                    for i in range(0, total_chunks, batch_size):
                        batch = text_chunks[i:i+batch_size]
                        batch_embeddings = model.encode(batch)
                        all_embeddings.extend(batch_embeddings)
                        progress = 0.5 + 0.4 * (min(i + batch_size, total_chunks) / total_chunks)
                        callback(f"已向量化 {min(i + batch_size, total_chunks)}/{total_chunks} 个文本块", progress=progress)

                    return np.array(all_embeddings), model
                except Exception as e:
                    if attempt < max_retries - 1:
                        wait_time = 2 ** attempt
                        callback(f"尝试加载模型失败，将在 {wait_time} 秒后重试... ({attempt+1}/{max_retries})")
                        sleep(wait_time)
                    else:
                        callback(f"无法加载模型，请检查网络连接。错误: {str(e)}")
                        raise

        vectors, model = embed_with_progress(chunks)

        callback("构建索引中...", progress=0.9)
        build_index(vectors)

        with open("chunks.txt", "w", encoding="utf-8") as f:
            f.write("\n\n===CHUNK===\n\n".join(chunks))

        callback("✅ 构建完成！已保存索引与文本块。", progress=1.0)

    def on_ask_click():
        query = question_entry.get()
        if not query.strip():
            messagebox.showwarning("提示", "请输入问题！")
            return

        ask_button.config(state=tk.DISABLED)
        answer_text.config(state=tk.NORMAL)
        answer_text.delete(1.0, tk.END)
        answer_text.insert(tk.END, "正在思考中...")
        answer_text.config(state=tk.DISABLED)

        def ask_thread():
            try:
                ask_question(query, display_answer)
            except Exception as e:
                display_answer(f"发生错误: {str(e)}")
            finally:
                ask_button.config(state=tk.NORMAL)

        threading.Thread(target=ask_thread).start()

    def on_upload_files():
        files = filedialog.askopenfilenames(
            title="选择文档文件",
            filetypes=[
                ("文档文件", "*.docx *.pptx *.pdf *.txt"),
                ("Word文档", "*.docx"),
                ("PowerPoint文档", "*.pptx"),
                ("PDF文档", "*.pdf"),
                ("文本文件", "*.txt"),
                ("所有文件", "*.*")
            ]
        )

        if not files:
            return

        if not os.path.exists("docs"):
            os.makedirs("docs")

        for file in files:
            import shutil
            filename = os.path.basename(file)
            dest = os.path.join("docs", filename)
            shutil.copy2(file, dest)
            update_status(f"已添加: {filename}")

    # 创建主窗口
    root = tk.Tk()
    root.title("智能文档问答系统")
    root.geometry("800x600")

    # 上传文件按钮
    upload_frame = tk.Frame(root)
    upload_frame.pack(fill=tk.X, padx=10, pady=5)

    upload_button = tk.Button(upload_frame, text="上传文档", command=on_upload_files)
    upload_button.pack(side=tk.LEFT)

    # 构建知识库按钮
    build_button = tk.Button(upload_frame, text="构建知识库", command=on_build_click)
    build_button.pack(side=tk.LEFT, padx=5)

    # 进度条框架
    progress_frame = tk.Frame(root)
    progress_frame.pack(fill=tk.X, padx=10, pady=2)

    # 进度变量和标签
    progress_var = tk.IntVar(value=0)
    progress_label = tk.Label(progress_frame, text="0%", width=5)
    progress_label.grid(row=0, column=1, padx=(0, 5))

    # 进度条
    progress_bar = ttk.Progressbar(progress_frame, variable=progress_var, maximum=100)
    # 初始时隐藏进度条
    progress_bar.grid_remove()

    # 状态面板
    status_frame = tk.LabelFrame(root, text="构建状态")
    status_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)

    status_text = scrolledtext.ScrolledText(status_frame, height=6)
    status_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
    status_text.insert(tk.END, "欢迎使用智能文档问答系统！\n请上传文档并构建知识库。\n")
    status_text.config(state=tk.DISABLED)

    # 问题输入框
    question_frame = tk.Frame(root)
    question_frame.pack(fill=tk.X, padx=10, pady=5)

    tk.Label(question_frame, text="问题:").pack(side=tk.LEFT)
    question_entry = tk.Entry(question_frame)
    question_entry.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=5)

    ask_button = tk.Button(question_frame, text="提问", command=on_ask_click)
    ask_button.pack(side=tk.LEFT)

    # 回答面板
    answer_frame = tk.LabelFrame(root, text="回答")
    answer_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)

    answer_text = scrolledtext.ScrolledText(answer_frame)
    answer_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
    answer_text.config(state=tk.DISABLED)

    # 开始主循环
    root.mainloop()

if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1 and sys.argv[1] == "--cli":
        cli_interface()
    else:
        gui_interface()